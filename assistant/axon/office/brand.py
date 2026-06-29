"""Brand / theme presets: store the user's colors, font, and logo once, then apply them across a
whole PowerPoint deck or Word document in one step ("brand it" / "make it prettier")."""
import os

from pptx import Presentation
from pptx.util import Inches
from pptx.dml.color import RGBColor

from docx import Document
from docx.shared import Pt as DocxPt, RGBColor as DocxRGB

from axon.util import _result
from axon.settings import load_settings, save_settings
from axon.office.preview import ensure_closed, open_doc, save_with_retry

DEFAULT_BRAND = {"primary": "1F3A5F", "accent": "E07A2F", "dark": "0B0B0E",
                 "light": "FFFFFF", "text": "222222", "font": "Calibri", "logo": ""}


def get_brand():
    b = dict(DEFAULT_BRAND)
    b.update(load_settings().get("brand") or {})
    return b


def set_brand(args):
    """Save brand settings (primary/accent/dark/light/text colors as RRGGBB, font name, logo path)."""
    s = load_settings()
    b = s.get("brand") or {}
    for k in ("primary", "accent", "dark", "light", "text", "font", "logo"):
        if args.get(k) is not None:
            b[k] = str(args[k]).lstrip("#") if k not in ("font", "logo") else str(args[k])
    s["brand"] = b
    save_settings(s)
    return _result(f"Brand saved: {get_brand()}")


def _c(hexstr):
    return RGBColor.from_string(str(hexstr).lstrip("#"))


def brand_deck(args):
    """Apply the saved brand to every slide of a deck: title slide gets the dark background + light
    title; content slides get the light background, brand-colored titles, and consistent body text;
    the logo (if set) is placed top-right. Use this to brand a deck or 'make it prettier'."""
    path = args.get("path")
    if not path or not os.path.exists(path):
        return _result(f"No such deck: {path}", True)
    b = get_brand()
    ensure_closed(path)
    try:
        prs = Presentation(path)
    except Exception as e:
        return _result(f"Could not open deck: {e}", True)
    for i, s in enumerate(prs.slides):
        is_title = (i == 0)
        s.background.fill.solid()
        s.background.fill.fore_color.rgb = _c(b["dark"] if is_title else b["light"])
        for shape in s.shapes:
            if not shape.has_text_frame:
                continue
            is_title_shape = (shape == s.shapes.title)
            for p in shape.text_frame.paragraphs:
                for r in p.runs:
                    r.font.name = b["font"]
                    if is_title_shape:
                        r.font.bold = True
                        r.font.color.rgb = _c(b["light"] if is_title else b["primary"])
                    else:
                        r.font.color.rgb = _c(b["light"] if is_title else b["text"])
        if b.get("logo") and os.path.exists(b["logo"]) and not is_title:
            try:
                s.shapes.add_picture(b["logo"], Inches(8.4), Inches(0.25), height=Inches(0.55))
            except Exception:
                pass
    final, err = save_with_retry(path, lambda p: prs.save(p))
    if err is not None:
        return _result(f"Could not save deck: {err}", True)
    open_doc(final)
    msg = f"Branded {len(prs.slides)} slides in {final}."
    if final != path:
        msg += " (Original was open/locked, saved a copy.)"
    return _result(msg)


def brand_doc(args):
    """Apply the saved brand to a Word document: brand-colored headings in the brand font, and a
    consistent body font."""
    path = args.get("path")
    if not path or not os.path.exists(path):
        return _result(f"No such document: {path}", True)
    b = get_brand()
    ensure_closed(path)
    try:
        doc = Document(path)
    except Exception as e:
        return _result(f"Could not open document: {e}", True)
    prim = DocxRGB.from_string(str(b["primary"]).lstrip("#"))
    txt = DocxRGB.from_string(str(b["text"]).lstrip("#"))
    for p in doc.paragraphs:
        st = p.style.name if p.style else ""
        is_head = st.startswith("Heading") or st == "Title"
        for r in p.runs:
            r.font.name = b["font"]
            r.font.color.rgb = prim if is_head else txt
            if is_head:
                r.font.bold = True
    final, err = save_with_retry(path, lambda p: doc.save(p))
    if err is not None:
        return _result(f"Could not save document: {err}", True)
    open_doc(final)
    msg = f"Branded {final}."
    if final != path:
        msg += " (Original was open/locked, saved a copy.)"
    return _result(msg)
