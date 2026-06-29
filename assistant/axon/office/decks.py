"""PowerPoint via python-pptx (structured editing): read a deck's outline and apply targeted edits
(add/move/delete slides, titles, bullets, images, tables, charts, notes, backgrounds). Saves the
file, then opens it in PowerPoint as a live preview. Creates the deck if the path doesn't exist.
"""
import os

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE

from axon.util import _result
from axon.office._base import _resolve_path
from axon.office.preview import ensure_closed, open_doc, save_with_retry

# Friendly layout names -> the standard template layout indexes.
_LAYOUTS = {"title": 0, "title_content": 1, "section": 2, "two_content": 3,
            "comparison": 4, "title_only": 5, "blank": 6}
_CHART_TYPES = {"column": XL_CHART_TYPE.COLUMN_CLUSTERED, "bar": XL_CHART_TYPE.BAR_CLUSTERED,
                "line": XL_CHART_TYPE.LINE, "pie": XL_CHART_TYPE.PIE,
                "area": XL_CHART_TYPE.AREA}


def _rgb(s):
    return RGBColor.from_string(str(s).lstrip("#"))


def _slide(prs, n):
    slides = prs.slides
    i = int(n) - 1
    if i < 0 or i >= len(slides):
        raise IndexError(f"slide {n} out of range (deck has {len(slides)} slides)")
    return slides[i]


def _body_placeholder(slide):
    """The slide's content/body placeholder (for bullets), or None."""
    for ph in slide.placeholders:
        if ph.placeholder_format.idx != 0 and ph.has_text_frame:  # idx 0 is the title
            return ph
    return None


def _set_bullets(slide, bullets):
    ph = _body_placeholder(slide)
    if ph is None:
        return "no body placeholder on this slide (use a title_content layout)"
    tf = ph.text_frame
    tf.clear()
    for i, b in enumerate(bullets or []):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = str(b)
    return None


def _add_slide(prs, op):
    layout_idx = _LAYOUTS.get(str(op.get("layout", "title_content")), 1)
    layout_idx = min(layout_idx, len(prs.slide_layouts) - 1)
    slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
    if op.get("title") and slide.shapes.title is not None:
        slide.shapes.title.text = str(op["title"])
    if op.get("bullets"):
        _set_bullets(slide, op["bullets"])
    if op.get("notes"):
        slide.notes_slide.notes_text_frame.text = str(op["notes"])
    at = op.get("at")
    if at is not None:  # move the newly-appended slide to position `at` (1-based)
        _move_slide(prs, len(prs.slides), int(at))
    return None


def _move_slide(prs, frm, to):
    lst = prs.slides._sldIdLst
    ids = list(lst)
    frm, to = int(frm) - 1, max(0, min(int(to) - 1, len(ids) - 1))
    el = ids[frm]
    lst.remove(el)
    lst.insert(to, el)
    return None


def _delete_slide(prs, n):
    lst = prs.slides._sldIdLst
    ids = list(lst)
    lst.remove(ids[int(n) - 1])
    return None


def _add_image(slide, op):
    path = op.get("path")
    if not path or not os.path.exists(path):
        return f"image not found: {path}"
    kw = {}
    if op.get("width"):
        kw["width"] = Inches(float(op["width"]))
    if op.get("height"):
        kw["height"] = Inches(float(op["height"]))
    slide.shapes.add_picture(path, Inches(float(op.get("left", 1))), Inches(float(op.get("top", 1))), **kw)
    return None


def _add_table(slide, op):
    rows = op.get("rows") or []
    if not rows:
        return "table needs rows"
    nrows, ncols = len(rows), max(len(r) for r in rows)
    gt = slide.shapes.add_table(nrows, ncols, Inches(float(op.get("left", 0.5))),
                                Inches(float(op.get("top", 1.5))), Inches(float(op.get("width", 9))),
                                Inches(float(op.get("height", 0.4 * nrows)))).table
    for r, row in enumerate(rows):
        for c in range(ncols):
            gt.cell(r, c).text = str(row[c]) if c < len(row) else ""
    return None


def _add_chart(slide, op):
    ct = _CHART_TYPES.get(str(op.get("chart_type", "column")), XL_CHART_TYPE.COLUMN_CLUSTERED)
    data = CategoryChartData()
    data.categories = op.get("categories") or []
    series = op.get("series") or {}
    for name, vals in series.items():
        data.add_series(str(name), tuple(vals))
    gframe = slide.shapes.add_chart(ct, Inches(float(op.get("left", 1))), Inches(float(op.get("top", 1.5))),
                                    Inches(float(op.get("width", 8))), Inches(float(op.get("height", 4.5))), data)
    if op.get("title"):
        gframe.chart.has_title = True
        gframe.chart.chart_title.text_frame.text = str(op["title"])
    return None


def _set_bg(prs, op):
    color = _rgb(op.get("color", "FFFFFF"))
    which = op.get("slide", "all")
    targets = prs.slides if which == "all" else [_slide(prs, which)]
    for s in targets:
        s.background.fill.solid()
        s.background.fill.fore_color.rgb = color
    return None


def _apply_op(prs, op):
    kind = (op.get("op") or "").lower()
    if kind == "add_slide":
        return _add_slide(prs, op)
    if kind == "delete_slide":
        return _delete_slide(prs, op["slide"])
    if kind == "move_slide":
        return _move_slide(prs, op["slide"], op["to"])
    if kind == "set_title":
        s = _slide(prs, op["slide"])
        if s.shapes.title is None:
            return "this slide has no title placeholder"
        s.shapes.title.text = str(op.get("text", ""))
        return None
    if kind == "set_bullets":
        return _set_bullets(_slide(prs, op["slide"]), op.get("bullets", []))
    if kind == "set_notes":
        _slide(prs, op["slide"]).notes_slide.notes_text_frame.text = str(op.get("text", ""))
        return None
    if kind == "add_image":
        return _add_image(_slide(prs, op["slide"]), op)
    if kind == "add_table":
        return _add_table(_slide(prs, op["slide"]), op)
    if kind == "add_chart":
        return _add_chart(_slide(prs, op["slide"]), op)
    if kind == "set_bg":
        return _set_bg(prs, op)
    return f"unknown op: {kind}"


def ppt_read(args):
    """Outline a presentation: per slide -> index, layout, title, bullet lines, other shapes, notes."""
    path = args.get("path")
    if not path or not os.path.exists(path):
        return _result(f"No such file: {path}", True)
    try:
        prs = Presentation(path)
    except Exception as e:
        return _result(f"Could not open deck: {e}", True)
    lines = [f"{os.path.basename(path)} — {len(prs.slides)} slides"]
    for i, s in enumerate(prs.slides, 1):
        title = s.shapes.title.text if s.shapes.title is not None and s.shapes.title.text else "(no title)"
        lines.append(f"\nSlide {i}: {title}")
        body = _body_placeholder(s)
        if body is not None and body.text_frame.text.strip():
            for p in body.text_frame.paragraphs:
                if p.text.strip():
                    lines.append(f"  • {p.text}")
        extras = [sh.shape_type and str(sh.shape_type).split()[0].lower() or sh.name
                  for sh in s.shapes if sh != s.shapes.title and not sh.is_placeholder]
        extras = [x for x in extras if x]
        if extras:
            lines.append(f"  [shapes: {', '.join(extras)}]")
        try:
            note = s.notes_slide.notes_text_frame.text.strip() if s.has_notes_slide else ""
        except Exception:
            note = ""
        if note:
            lines.append(f"  [notes: {note[:160]}]")
    return _result("\n".join(lines))


def ppt_edit(args):
    """Apply a list of edit ops to a deck (creating it if missing), save, and open it in PowerPoint.

    ops example:
      [{"op":"add_slide","layout":"title","title":"Q3 Review","bullets":["Revenue","Risks"]},
       {"op":"add_image","slide":1,"path":"C:/logo.png","left":8,"top":0.3,"width":1.2},
       {"op":"add_chart","slide":2,"chart_type":"column","categories":["Jan","Feb"],
        "series":{"Sales":[10,14]},"title":"Sales"}]
    """
    path = _resolve_path(args.get("path"), "presentation", ".pptx")
    ops = args.get("ops") or []
    if not ops:
        return _result("Provide ops: a list of edit operations (see add_slide / set_bullets / add_chart ...).", True)
    ensure_closed(path)  # unlock the file if it's open in PowerPoint
    try:
        prs = Presentation(path) if os.path.exists(path) else Presentation()
    except Exception as e:
        return _result(f"Could not open deck: {e}", True)
    done, errors = 0, []
    for n, op in enumerate(ops, 1):
        try:
            err = _apply_op(prs, op)
            if err:
                errors.append(f"op {n} ({op.get('op')}): {err}")
            else:
                done += 1
        except Exception as e:
            errors.append(f"op {n} ({op.get('op')}): {e}")
    final, serr = save_with_retry(path, lambda p: prs.save(p))
    if serr is not None:
        return _result(f"Could not save deck: {serr}", True)
    open_doc(final)  # live preview on the dot's monitor
    msg = f"Applied {done}/{len(ops)} edits to {final} ({len(prs.slides)} slides)."
    if final != path:
        msg += " (Original was open/locked, so I saved an updated copy.)"
    if errors:
        msg += " Issues:\n" + "\n".join(errors)
    return _result(msg, bool(errors) and done == 0)
