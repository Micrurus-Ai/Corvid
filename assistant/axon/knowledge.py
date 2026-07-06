"""Project/folder context + document Q&A across the user's files.

RAG-lite: read supported files in a folder, keyword-rank the most relevant chunks, and let the LLM
compose an answer with [filename] citations. No embeddings/index needed — reliable and dependency-light.
"""
import os
import re

from openai import OpenAI

from axon.util import _result
from axon.config import MODEL
from axon.settings import load_settings, save_settings

_EXT = {".txt", ".md", ".csv", ".log", ".docx", ".pdf", ".pptx", ".xlsx",
        ".htm", ".html", ".json"}


def set_project(args):
    """Set the active project folder so 'ask my documents' searches it by default."""
    folder = args.get("folder")
    if not folder or not os.path.isdir(folder):
        return _result(f"Not a folder: {folder}", True)
    s = load_settings()
    s["project_folder"] = folder
    save_settings(s)
    return _result(f"Project folder set: {folder}")


def _read_text(path):
    ext = os.path.splitext(path)[1].lower()
    try:
        if ext in (".txt", ".md", ".csv", ".log", ".json"):
            with open(path, encoding="utf-8", errors="ignore") as f:
                return f.read()
        if ext in (".htm", ".html"):
            with open(path, encoding="utf-8", errors="ignore") as f:
                raw = f.read()
            raw = re.sub(r"(?is)<(script|style)[^>]*>.*?</\1>", " ", raw)  # drop scripts/styles
            return re.sub(r"<[^>]+>", " ", raw)                            # strip remaining tags
        if ext == ".docx":
            from docx import Document
            return "\n".join(p.text for p in Document(path).paragraphs)
        if ext == ".pdf":
            from pypdf import PdfReader
            return "\n".join((pg.extract_text() or "") for pg in PdfReader(path).pages)
        if ext == ".pptx":
            from pptx import Presentation
            out = []
            for s in Presentation(path).slides:
                for sh in s.shapes:
                    if sh.has_text_frame:
                        out.append(sh.text_frame.text)
            return "\n".join(out)
        if ext == ".xlsx":
            from openpyxl import load_workbook
            wb = load_workbook(path, data_only=True, read_only=True)
            out = []
            for ws in wb.worksheets:
                for row in ws.iter_rows(values_only=True):
                    out.append(" ".join("" if c is None else str(c) for c in row))
            return "\n".join(out)
    except Exception:
        return ""
    return ""


def _chunks(text, path, size=1200):
    out = []
    for i in range(0, len(text), size):
        c = text[i:i + size].strip()
        if c:
            out.append((path, c))
    return out


def _keyword_top(folder, question):
    """Fallback retrieval: read files, keyword-rank chunks (caps at ~400 files). Returns [(path, chunk)]."""
    chunks = []
    scanned = 0
    for root, _dirs, files in os.walk(folder):
        for fn in files:
            if os.path.splitext(fn)[1].lower() in _EXT:
                txt = _read_text(os.path.join(root, fn))
                if txt:
                    chunks.extend(_chunks(txt, os.path.join(root, fn))[:30])
                    scanned += 1
        if scanned > 400:
            break
    if not chunks:
        return []
    terms = [w for w in re.findall(r"\w+", question.lower()) if len(w) > 2]

    def score(c):
        t = c[1].lower()
        return sum(t.count(w) for w in terms)

    ranked = sorted(chunks, key=score, reverse=True)
    return [c for c in ranked if score(c) > 0][:8] or ranked[:4]


def ask_documents(args):
    """Answer a question using the user's documents (args: question, folder?). Searches the project
    folder if no folder is given. Uses a semantic embeddings index (fast + no file cap on big folders),
    falling back to a keyword scan. Returns an answer with [filename] citations."""
    question = args.get("question")
    folder = args.get("folder") or load_settings().get("project_folder")
    if not question:
        return _result("Provide a question.", True)
    if not folder or not os.path.isdir(folder):
        return _result("No project folder set — use set_project, or pass a folder.", True)
    if not os.getenv("OPENAI_API_KEY"):
        return _result("No API key available.", True)

    top = []
    try:
        from axon.doc_index import query
        top = query(folder, question, k=10)   # semantic, indexed, incremental
    except Exception:
        top = []
    if not top:
        top = _keyword_top(folder, question)   # fallback
    if not top:
        return _result(f"No readable documents found in {folder}.", True)
    context = "\n\n".join(f"[{os.path.basename(p)}]\n{c}" for p, c in top)
    try:
        from axon.llm import text_llm
        client, model = text_llm()
        prompt = (
            "Answer the question using ONLY the document excerpts below. Cite the source file name(s) "
            "in [brackets] for each claim. If the answer isn't in the excerpts, say so plainly.\n\n"
            f"EXCERPTS:\n{context[:12000]}\n\nQUESTION: {question}")
        resp = client.chat.completions.create(
            model=model, messages=[{"role": "user", "content": prompt}], temperature=0)
        ans = resp.choices[0].message.content or ""
    except Exception as e:
        return _result(f"Q&A failed: {e}", True)
    scanned_srcs = sorted(set(os.path.basename(p) for p, _ in top))
    cited = sorted(set(re.findall(r"\[([^\[\]]+\.[A-Za-z0-9]{1,5})\]", ans)))
    footer = "\n\nSources cited: " + ", ".join(cited) if cited else \
             "\n\nSources scanned: " + ", ".join(scanned_srcs)
    return _result(ans + footer)


def index_documents(args):
    """Build/refresh the semantic index for a folder so Q&A over large folders is fast and complete."""
    folder = args.get("folder") or load_settings().get("project_folder")
    if not folder or not os.path.isdir(folder):
        return _result("No folder to index — set a project folder or pass a folder.", True)
    if not os.getenv("OPENAI_API_KEY"):
        return _result("No API key available.", True)
    try:
        from axon.doc_index import index_stats
        st = index_stats(folder)
    except Exception as e:
        return _result(f"Indexing failed: {e}", True)
    if not st:
        return _result(f"No readable documents found in {folder}.", True)
    return _result(f"Indexed {st['files']} file(s) into {st['chunks']} chunk(s) from {folder}. "
                   "Q&A now covers the whole folder and is fast on repeat questions.")
